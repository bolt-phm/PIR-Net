import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 资源目录
IMG_DIR = "ppt_images"

# --- 占位符图片生成器 (辅助) ---
# 如果您没有把论文原图放进去，这个函数会生成一个临时的红色提示图
def create_missing_img_placeholder(filename, text_hint):
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots(figsize=(6, 4))
    ax.text(0.5, 0.5, f"【请在此处插入原图】\n{filename}\n\n{text_hint}", 
            ha='center', va='center', fontsize=15, color='red', fontweight='bold')
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_edgecolor('red')
        spine.set_linewidth(5)
        spine.set_linestyle('--')
    path = os.path.join(IMG_DIR, "MISSING_" + filename)
    fig.savefig(path, bbox_inches='tight')
    plt.close(fig)
    return path

class PIRNetPPTBuilder:
    def __init__(self):
        self.prs = Presentation()
        # 16:9 宽屏
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.layout_blank = self.prs.slide_layouts[6]

    def _get_image(self, filename, hint=""):
        """智能获取图片：优先找生成的数据图 -> 其次找用户放入的原图 -> 最后生成占位符"""
        path = os.path.join(IMG_DIR, filename)
        
        # 1. 检查是否存在
        if os.path.exists(path):
            return path
        
        # 2. 如果不存在，生成一个显眼的占位图
        print(f"提示：未找到图片 {filename}，生成占位符供您后续替换。")
        return create_missing_img_placeholder(filename, hint)

    def add_slide(self, title_text, content_points, img_filename=None, img_hint="", layout="split"):
        """
        layout: 'split' (左文右图), 'full_bottom' (上文下图), 'center' (纯文字)
        """
        slide = self.prs.slides.add_slide(self.layout_blank)
        
        # 1. 标题栏
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = title_box.text_frame
        tf.text = title_text
        p = tf.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "微软雅黑"
        p.font.color.rgb = RGBColor(0, 51, 102) # 深海军蓝

        # 2. 内容布局
        if layout == "split":
            # 左侧文本
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5))
            # 右侧图片位置
            img_left, img_top, img_w = Inches(7.2), Inches(1.8), Inches(5.8)
        elif layout == "full_bottom":
            # 顶部文本
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0))
            # 底部图片位置
            img_left, img_top, img_h = Inches(1.5), Inches(3.6), Inches(3.5)
        else:
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        
        # 3. 填充文本内容
        tf_content = text_box.text_frame
        tf_content.word_wrap = True
        for point in content_points:
            p = tf_content.add_paragraph()
            p.text = point
            p.font.name = "微软雅黑"
            p.font.size = Pt(22) # 增大字号
            p.space_after = Pt(18)
            p.level = 0
            if point.startswith(">>"): # 二级标题逻辑
                p.level = 1
                p.text = point.replace(">>", "")
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(80, 80, 80)

        # 4. 插入图片
        if img_filename:
            img_path = self._get_image(img_filename, img_hint)
            if layout == "full_bottom":
                slide.shapes.add_picture(img_path, img_left, img_top, height=img_h)
            elif layout == "split":
                # 保持比例缩放
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_w)

    def build(self):
        # --- 封面 ---
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "PIR-Net: 面向超高频螺栓松动检测的\n物理感知重采样与非对称融合网络"
        subtitle.text = "解决 1MHz 高频采样下的“效率与特征保留”矛盾\n\n汇报人：XXX"
        title.text_frame.paragraphs[0].font.name = "微软雅黑"
        title.text_frame.paragraphs[0].font.bold = True
        
        # --- Slide 1: 背景与痛点 ---
        self.add_slide(
            "研究背景与核心矛盾",
            [
                "工业背景：螺栓松动是桥梁、风机等结构失效的主要原因。",
                "现有技术的局限性：",
                ">> 视觉方法：无法检测内部应力丢失，受光照影响大。",
                ">> 传统振动：低频采样无法捕捉微秒级冲击 (Micro-transients)。",
                "核心矛盾 (The Dilemma)：",
                ">> 要捕捉 1-10µs 的微弱冲击，必须使用 1MHz 超高频采样。",
                ">> 边缘端算力受限，无法实时处理海量数据。",
                ">> 传统降采样 (Downsampling) 会直接抹除冲击特征（混叠效应）。"
            ],
            "resampling_comparison_detailed.png", # 使用刚才代码生成的对比图
            "此处展示：线性降采样导致冲击丢失 vs PIR保留冲击"
        )

        # --- Slide 2: PIR-Net 解决方案总览 ---
        self.add_slide(
            "PIR-Net 整体架构设计",
            [
                "设计理念：物理感知 (Physics-Informed) + 非对称融合。",
                "核心组件 (2-2-2 架构)：",
                "1. 物理感知自适应重采样 (150:1 压缩率)",
                "2. 5通道异构张量表示 (梯度+能量+时域)",
                "3. 非对称多头注意力融合 (Asymmetric Fusion)",
                "优势：在边缘端实现 36ms 实时推理，精度达 96.04%。"
            ],
            "moxingzonglan.png", # 【注意】这里需要您放入论文原图
            "请插入论文 Figure 3: Model Architecture",
            layout="full_bottom"
        )

        # --- Slide 3: 创新点 I - 物理感知重采样 ---
        self.add_slide(
            "创新点 I：物理感知重采样机制",
            [
                "传统方法的失效：线性采样将稀疏冲击视为噪声并平均化。",
                "PIR 算法核心：",
                ">> Max-Pooling (0.7): 强行保留微秒级冲击峰值 (Impact)。",
                ">> Mean-Pooling (0.3): 兼顾整体能量趋势 (Energy)。",
                "效果：",
                "在 150倍 压缩下，信号信噪比 (SNR) 不降反升。",
                "完美保留了早期松动的微弱故障特征。"
            ],
            "preprocessing_demo.png", # 建议放入论文 Figure 2
            "请插入论文 Figure 2: Resampling Comparison",
            layout="split"
        )

        # --- Slide 4: 创新点 II - 5通道异构特征 ---
        self.add_slide(
            "创新点 II：5通道异构梯度特征",
            [
                "痛点：STFT 频谱图丢失了相位和波形形态信息。",
                "解决方案：构建 5-Channel Tensor",
                "1. Log频谱 (频域强度)",
                "2. 时间梯度 (捕捉冲击时刻)",
                "3. 频率梯度 (捕捉共振漂移)",
                "4. 能量图 (聚焦高能区)",
                "5. 时域嵌入 (原始波形形态)",
                "贡献：梯度特征显著提升了“过渡态”识别率。"
            ],
            "5channel_feature.png", # 建议放入论文 Figure 4
            "请插入论文 Figure 4: 5-Channel Representation",
            layout="split"
        )

        # --- Slide 5: 创新点 III - 非对称融合 ---
        self.add_slide(
            "创新点 III：非对称注意力融合",
            [
                "策略：“信号主导，图像补偿” (Signal-Dominant)",
                "物理依据：",
                ">> 1MHz 原始信号的信息密度远高于 224x224 图像。",
                ">> 图像分支不应喧宾夺主，而是作为“安全冗余”。",
                "实现方式：",
                "多头注意力机制 (Multi-Head Attention)",
                "让稳定的频谱特征去“检索”高噪的时域细节。"
            ],
            "moxingliuchengtu.png", # 建议放入论文 Figure 1 或 流程图
            "请插入论文流程图或 Attention 示意图",
            layout="split"
        )

        # --- Slide 6: 实验结果 (SOTA) ---
        self.add_slide(
            "主实验结果：SOTA 性能突破",
            [
                "综合准确率：96.04% (提升 +5.32%)",
                "关键指标突破：",
                ">> 过渡态 (Transition) F1-Score: 98.62%",
                ">> 严重松动 (Severe Loose) 漏检率: 0%",
                "鲁棒性：",
                "在 20dB 噪声下仅下降 1.17%，远优于 CNN-LSTM。",
                "推理速度：36.2ms (满足工业实时性要求)。"
            ],
            "chart_sota.png", # 使用代码生成的柱状图
            layout="split"
        )

        # --- Slide 7: 消融实验与误差分析 ---
        self.add_slide(
            "深入分析：消融与误差来源",
            [
                "消融实验结论：",
                ">> 多头注意力机制贡献最大 (+5.32%)。",
                ">> 单纯拼接 (Concat) 无法挖掘模态互补性。",
                "误差分析 (Confusion Matrix)：",
                ">> “过渡态”识别非常精准 (误差 < 1.6%)。",
                ">> 真正难点：区分“紧固 (Secure)”与“过紧 (Over-tight)”。",
                ">> 原因：两者在物理刚度上接近，非线性特征差异微小。"
            ],
            "confusion_matrix_comparison_fixed.png", # 建议插入论文 Figure 7
            "请插入论文 Figure 7: Confusion Matrix",
            layout="split"
        )

        # --- Slide 8: 可视化分析 (t-SNE) ---
        self.add_slide(
            "特征空间可视化 (t-SNE)",
            [
                "聚类效果：",
                "六种松动状态在特征空间内清晰分离。",
                "拓扑逻辑正确：",
                ">> “过渡态”集群正确地位于“松动”与“紧固”之间。",
                "安全边界 (Safety Margin)：",
                ">> 严重松动 (蓝色) 与其他类别距离极远。",
                ">> 解释了为何能实现“零漏检”。"
            ],
            "fig5_tsne.png", # 建议插入论文 Figure 5
            "请插入论文 Figure 5: t-SNE Scatter Plot",
            layout="full_bottom"
        )

        # --- Slide 9: 核心机理 (高光时刻) ---
        self.add_slide(
            "高光时刻：动态补偿机制",
            [
                "问题：为什么要引入图像模态？",
                "发现：Outlier Analysis (离群点分析)",
                ">> 平均情况下：信号权重 > 99% (主导)。",
                ">> 困难样本下：图像权重激增至 48%。",
                "机理解释：",
                "当信号特征模糊时（如紧固/过紧边界），模型自动",
                "激活视觉通路，利用纹理特征进行修正。",
                "结论：图像模态充当了系统的“安全气囊”。"
            ],
            "chart_mechanism.png", # 使用代码生成的高光图表
            layout="split"
        )

        # --- Slide 10: 总结 ---
        self.add_slide(
            "结论与展望",
            [
                "总结：",
                "PIR-Net 成功解决了超高频 PHM 的效率瓶颈。",
                "提出了“物理感知重采样”与“非对称融合”新范式。",
                "核心贡献：",
                "1. 实现了 150:1 的高保真数据压缩。",
                "2. 解决了“过渡态”识别难题 (F1 98.6%)。",
                "3. 建立了动态互补的安全性机制。",
                "未来工作：推广至轴承故障与齿轮箱监测。"
            ],
            None,
            layout="center"
        )

        output_file = "PIR_Net_汇报胶片_CN.pptx"
        self.prs.save(output_file)
        print(f"PPT 生成成功！文件名为: {output_file}")
        print("请打开 PPT，将带有红色边框的占位图替换为您的论文原图。")

if __name__ == "__main__":
    builder = PIRNetPPTBuilder()
    builder.build()